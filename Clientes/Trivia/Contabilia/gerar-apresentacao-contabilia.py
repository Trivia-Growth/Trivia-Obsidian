#!/usr/bin/env python3
"""
Gera a apresentação PowerPoint da plataforma Contabilia.
Executar: python3 gerar-apresentacao-contabilia.py
Resultado: apresentacao-contabilia.pptx no mesmo diretório do script.
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    import subprocess
    print("python-pptx não encontrado. Instalando...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Cores do tema ---
BG_COLOR = RGBColor(15, 23, 42)        # slate-900
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240)   # slate-200
ACCENT_BLUE = RGBColor(59, 130, 246)   # blue-500
ACCENT_GREEN = RGBColor(16, 185, 129)  # green-500


def set_slide_bg(slide, color):
    """Define cor de fundo sólida para o slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, font_size=Pt(34), bold=True, color=WHITE, left=Inches(0.8), top=Inches(0.5), width=Inches(11.5), height=Inches(1.2)):
    """Adiciona título ao slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_bullets(slide, bullets, font_size=Pt(18), color=LIGHT_GRAY, left=Inches(1.0), top=Inches(1.8), width=Inches(11.0), height=Inches(5.0)):
    """Adiciona lista de bullets ao slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = font_size
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.level = 0
    return txBox


def add_footer(slide, text, color=LIGHT_GRAY):
    """Adiciona rodapé ao slide."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER


def create_presentation():
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Layout em branco
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1 — Capa
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)

    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Contabilia"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.3), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Automação inteligente para escritórios contábeis"
    p2.font.size = Pt(24)
    p2.font.bold = False
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    add_footer(slide, "Trivia Studio | 2026")

    # =========================================================================
    # SLIDE 2 — O Problema
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "A Tempestade Perfeita")
    add_bullets(slide, [
        "Reforma tributária (2026-2032): trabalho dobra com apuração dual",
        "Split payment fragmenta recebimentos e complica conciliação",
        "40-60% mais trabalho operacional, sem aumento de receita",
        '"Sou pago como operacional, cobrado como estratégico"',
    ])

    # =========================================================================
    # SLIDE 3 — Dores do Dia a Dia
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "O Que Tira Seu Sono")
    add_bullets(slide, [
        "Clientes atrasam documentos → prazo fiscal estoura",
        "Retrabalho por informação errada",
        'WhatsApp como "sistema de gestão"',
        "Inadimplência silenciosa de honorários",
        "Zero tempo para ser consultivo",
    ])

    # =========================================================================
    # SLIDE 4 — A Solução
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Contabilia: Sua Operação no Piloto Automático")
    add_bullets(slide, [
        "Plataforma SaaS que conecta escritório e cliente",
        "Automatiza: coleta → emissão → classificação → apuração → entrega",
        "IA que aprende o padrão do seu escritório",
        "Simples de usar (cliente e contador)",
    ])

    # =========================================================================
    # SLIDE 5 — CRM Contábil
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Base de Clientes Completa")
    add_bullets(slide, [
        "Cadastro com todos os dados fiscais",
        "Controle de honorários + cobrança automática",
        "Score de saúde: pontualidade, inadimplência, complexidade",
        "Régua de cobrança: D-3, D+0, D+3, D+7, D+15",
        "Visão de rentabilidade por cliente",
    ])

    # =========================================================================
    # SLIDE 6 — Emissão de Notas
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Emissão Inteligente de NFS-e")
    add_bullets(slide, [
        "Individual, em lote (grid ou Excel), ou recorrente",
        "Motor de tributos automático (ISS + IBS + CBS)",
        "IA sugere código de serviço, CFOP, NCM",
        "Validação pré-envio + cálculo em tempo real",
        "O CLIENTE pode emitir pelo portal (com supervisão)",
    ])

    # =========================================================================
    # SLIDE 7 — Emissão em Lote
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "12 Notas em 2 Minutos")
    add_bullets(slide, [
        "Importação de Excel com modelo pronto",
        "Grid tipo planilha para preencher online",
        "IA calcula tributos automaticamente por linha",
        "Preview completo antes de emitir",
        "Status em tempo real + download em lote (ZIP)",
    ])

    # =========================================================================
    # SLIDE 8 — Portal do Cliente
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "O Cliente Envia Organizado")
    add_bullets(slide, [
        'Checklist dinâmico: "o que falta enviar este mês"',
        "Upload por portal, WhatsApp ou email",
        "OCR + IA extrai dados de fotos e PDFs",
        "Cobrança automática de pendências (D-10, D-5, D-1)",
        "Zero WhatsApp pessoal para cobrar documento",
    ])

    # =========================================================================
    # SLIDE 9 — Atendimento
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Demandas Organizadas, SLA Cumprido")
    add_bullets(slide, [
        "Cliente abre chamado: portal ou WhatsApp",
        "IA classifica categoria e urgência automaticamente",
        "Fila com SLA por tipo de demanda",
        "Templates + sugestão de resposta por IA",
        "Pesquisa de satisfação ao concluir",
    ])

    # =========================================================================
    # SLIDE 10 — Comunicação
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "WhatsApp Profissional, Não Pessoal")
    add_bullets(slide, [
        "Régua automática de cobrança de documentos",
        "Envio de guias com vencimento",
        "Notificações de notas emitidas",
        "Resumo mensal automático para o cliente",
        "Histórico centralizado (nada se perde)",
    ])

    # =========================================================================
    # SLIDE 11 — Relatórios
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Dados que Geram Decisão")
    add_bullets(slide, [
        "Resumo fiscal mensal por cliente (enviado automaticamente)",
        "Comparativo regime antigo × novo (transição)",
        "Rentabilidade por cliente (custo real × honorário)",
        "Inadimplência, obrigações, demandas/SLA",
        "White-label: logo e marca do escritório",
    ])

    # =========================================================================
    # SLIDE 12 — IA na Prática
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "IA que Aprende com Você")
    add_bullets(slide, [
        "Classificação fiscal automática (95%+ acerto)",
        "Sugestão de código de serviço por descrição",
        "OCR de documentos (foto → dados estruturados)",
        "Detecção de anomalias fiscais",
        "Assistente tributário (legislação via RAG)",
        "Quanto mais usa, mais precisa fica",
    ])

    # =========================================================================
    # SLIDE 13 — Motor de Tributos
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Cálculo de Tributos 100% Automático")
    add_bullets(slide, [
        "Input: prestador + tomador + serviço + valor",
        "Output: ISS + IBS + CBS + retenções = valor líquido",
        "Regime detectado automaticamente",
        "Município destino → alíquota correta",
        "Transição (2026-2032): ambos os regimes",
        "Fontes: IBPT, BrasilAPI, tabelas municipais",
    ])

    # =========================================================================
    # SLIDE 14 — O Que o Cliente Vê
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Experiência Simples para a PME")
    add_bullets(slide, [
        "Enviar documentos (upload ou WhatsApp)",
        "Emitir notas (individual ou lote)",
        "Ver tributos calculados automaticamente",
        "Acompanhar guias e vencimentos",
        "Abrir demandas e acompanhar status",
        "Receber relatório fiscal mensal",
    ])

    # =========================================================================
    # SLIDE 15 — Impacto
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Resultados Esperados")
    add_bullets(slide, [
        "70% menos tempo em classificação fiscal",
        "Zero cobrança manual de documentos",
        "Inadimplência controlada com régua automática",
        "De 80% operacional → 80% consultivo",
        "ROI: 28h/semana liberadas × R$300/h = R$33.600/mês potencial",
    ])

    # =========================================================================
    # SLIDE 16 — Próximos Passos
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)
    add_title(slide, "Vamos Construir Juntos")
    add_bullets(slide, [
        "MVP em 12-16 semanas",
        "Stack: Supabase + Next.js + Nuvem Fiscal + Claude AI + Evolution API",
        "Modelo SaaS: a partir de R$ 297/mês",
        "Validação com escritório piloto",
        "Roadmap: NF-e, apuração dual, SPED, conciliação",
    ])
    add_footer(slide, "Trivia Studio — Automação Inteligente")

    # =========================================================================
    # Salvar arquivo
    # =========================================================================
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(script_dir, "apresentacao-contabilia.pptx")
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = create_presentation()
    print(f"Apresentação gerada com sucesso: {path}")
